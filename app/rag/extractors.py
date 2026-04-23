"""
app/rag/extractors.py
---------------------
Módulo de extracción de texto a partir de diferentes tipos de archivo.
Soporta: PDF, DOCX, XLSX, PPTX, HTML/web URL, texto plano.
Imágenes, Audio y Video quedan para Fase 2 (Google Vision / Speech-to-Text).
"""
import io
import re
import requests
from typing import Optional


def extract_from_pdf(file_bytes: bytes) -> str:
    """Extrae texto de un PDF usando pdfplumber (mejor con tablas y layouts complejos)."""
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        return "\n\n".join(text_parts)
    except Exception:
        # Fallback a PyPDF2
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def extract_from_docx(file_bytes: bytes) -> str:
    """Extrae texto de un archivo Word (.docx)."""
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # También extraer texto de tablas
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    return "\n\n".join(paragraphs)


def extract_from_xlsx(file_bytes: bytes) -> str:
    """Extrae datos de un archivo Excel (.xlsx) como texto tabular."""
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    text_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        text_parts.append(f"=== Hoja: {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                text_parts.append(row_text)
    return "\n".join(text_parts)


def extract_from_pptx(file_bytes: bytes) -> str:
    """Extrae texto de presentaciones PowerPoint (.pptx)."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    text_parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            text_parts.append(f"--- Diapositiva {i} ---\n" + "\n".join(slide_texts))
    return "\n\n".join(text_parts)


def extract_from_url(url: str) -> str:
    """Descarga y extrae texto de una página web."""
    from bs4 import BeautifulSoup
    headers = {"User-Agent": "Mozilla/5.0 (GammIA RAG Bot)"}
    response = requests.get(url, headers=headers, timeout=15)
    response.raise_for_status()
    soup = BeautifulSoup(response.text, "html.parser")
    # Eliminar scripts, estilos y nav
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    # Limpiar líneas vacías excesivas
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return "\n".join(lines)


def extract_from_txt(file_bytes: bytes) -> str:
    """Extrae texto plano de archivos .txt o .md."""
    try:
        return file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return file_bytes.decode("latin-1")


SUPPORTED_EXTENSIONS = {
    "pdf": extract_from_pdf,
    "docx": extract_from_docx,
    "doc": extract_from_docx,
    "xlsx": extract_from_xlsx,
    "xls": extract_from_xlsx,
    "pptx": extract_from_pptx,
    "ppt": extract_from_pptx,
    "txt": extract_from_txt,
    "md": extract_from_txt,
}


def extract_text_from_file(file_bytes: bytes, filename: str) -> Optional[str]:
    """
    Dispatcher principal. Extrae texto según la extensión del archivo.
    Retorna None si el tipo no está soportado.
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    extractor = SUPPORTED_EXTENSIONS.get(ext)
    if not extractor:
        return None
    text = extractor(file_bytes)
    # Limpiar caracteres no imprimibles
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
    return text.strip() or None
